import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from models.paper_structure import ResearchPaper

class PresentationGeneratorService:
    """Service for generating PowerPoint presentations from research papers"""
    
    def generate_presentation(self, paper: ResearchPaper, output_path: str):
        """Generate a PPTX file from the research paper"""
        prs = Presentation()
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = paper.title
        
        # Format authors for subtitle
        authors_str = ", ".join([a.name for a in paper.authors])
        subtitle.text = f"{authors_str}\n{paper.authors[0].affiliation}"
        
        # Content Slides
        self._add_content_slide(prs, "Abstract", paper.abstract)
        
        # Add sections
        section_order = ['introduction', 'methodology', 'results', 'discussion', 'conclusion']
        for section_name in section_order:
            if section_name in paper.sections and paper.sections[section_name]:
                title_text = section_name.replace('_', ' ').title()
                content = paper.sections[section_name]
                # Split long content into multiple slides
                self._add_content_slides_split(prs, title_text, content)
                
        # References Slide
        if paper.references:
            ref_text = ""
            for i, ref in enumerate(paper.references[:5], 1): # Limit to top 5 for slide
                ref_text += f"[{i}] {ref.title} ({ref.year})\n"
            self._add_content_slide(prs, "Key References", ref_text)
            
        prs.save(output_path)
        return output_path

    def _add_content_slide(self, prs, title_text, content_text):
        """Add a standard content slide"""
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title_text
        tf = body_shape.text_frame
        tf.text = content_text[:800] + "..." if len(content_text) > 800 else content_text
        
        # Basic formatting
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)

    def _add_content_slides_split(self, prs, title_text, content_text):
        """Add content slides, splitting text if too long"""
        # Simple splitting by sentences/length (approximate)
        max_chars = 600
        if len(content_text) <= max_chars:
            self._add_content_slide(prs, title_text, content_text)
            return
            
        parts = []
        current_part = ""
        sentences = content_text.split('. ')
        
        for sentence in sentences:
            if len(current_part) + len(sentence) < max_chars:
                current_part += sentence + ". "
            else:
                parts.append(current_part)
                current_part = sentence + ". "
        if current_part:
            parts.append(current_part)
            
        for i, part in enumerate(parts):
            slide_title = f"{title_text} ({i+1}/{len(parts)})" if len(parts) > 1 else title_text
            self._add_content_slide(prs, slide_title, part)
